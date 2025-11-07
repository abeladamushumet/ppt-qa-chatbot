import os
from pptx import Presentation
from ..utils.logger import logger
from ..config.settings import EXTRACTED_TEXT_DIR

# Ensure extracted_texts folder exists
os.makedirs(EXTRACTED_TEXT_DIR, exist_ok=True)


def extract_text_from_ppt(file_path: str) -> str:
    """
    Extract all text from a PowerPoint file.

    Args:
        file_path (str): Path to the PPTX file.

    Returns:
        str: Concatenated text from all slides.
    """
    try:
        prs = Presentation(file_path)
        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())

        full_text = "\n".join(text_runs)
        logger.info(f"Extracted text from PPT: {os.path.basename(file_path)}")
        return full_text

    except Exception as e:
        logger.error(f"Failed to extract text from {file_path}: {e}")
        return ""


def save_extracted_text(file_path: str, text: str) -> str:
    """
    Save the extracted text to the extracted_texts directory.

    Args:
        file_path (str): Original PPT file path.
        text (str): Extracted text.

    Returns:
        str: Path of the saved text file.
    """
    try:
        filename = os.path.basename(file_path)
        txt_filename = os.path.splitext(filename)[0] + ".txt"
        txt_path = os.path.join(EXTRACTED_TEXT_DIR, txt_filename)

        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(text)

        logger.info(f"Saved extracted text: {txt_filename}")
        return txt_path

    except Exception as e:
        logger.error(f"Failed to save extracted text for {file_path}: {e}")
        return ""


def process_ppt(file_path: str) -> str:
    """
    Complete pipeline: extract text and save to extracted_texts.
    Returns the path of the saved text file.
    """
    text = extract_text_from_ppt(file_path)
    if text:
        return save_extracted_text(file_path, text)
    else:
        logger.warning(f"No text extracted from {file_path}")
        return ""

if __name__ == "__main__":
    from ..config.settings import RAW_PPT_DIR

    for ppt_file in os.listdir(RAW_PPT_DIR):
        if ppt_file.endswith(".pptx"):
            ppt_path = os.path.join(RAW_PPT_DIR, ppt_file)
            txt_path = process_ppt(ppt_path)
            print(f"Saved extracted text to: {txt_path}")
